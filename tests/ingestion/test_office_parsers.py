"""Tests for DOCX and PPTX parsers."""

from __future__ import annotations

from pathlib import Path

import pytest

from custom_rag.core.types import BlockType
from custom_rag.ingestion.pipeline import parse_file

FIXTURES = Path(__file__).resolve().parents[1] / "fixtures"


@pytest.fixture(scope="module")
def sample_docx(tmp_path_factory: pytest.TempPathFactory) -> Path:
    pytest.importorskip("docx")
    from docx import Document

    path = tmp_path_factory.mktemp("office") / "sample.docx"
    document = Document()
    document.add_heading("Scope", level=1)
    document.add_paragraph("This document defines project scope.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Owner"
    table.rows[0].cells[1].text = "Platform"
    document.save(path)
    return path


@pytest.fixture(scope="module")
def sample_pptx(tmp_path_factory: pytest.TempPathFactory) -> Path:
    pytest.importorskip("pptx")
    from pptx import Presentation

    path = tmp_path_factory.mktemp("office") / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Update"
    slide.placeholders[1].text = "Revenue increased 12%"
    presentation.save(path)
    return path


def test_docx_parser_builds_section_and_table_rows(sample_docx: Path) -> None:
    document = parse_file(sample_docx)

    assert document.metadata.doc_type == "docx"
    sections = [block for block in document.blocks if block.block_type == BlockType.SECTION]
    assert sections
    rows = [block for block in document.blocks if block.block_type == BlockType.TABLE_ROW]
    assert rows
    assert "Owner" in rows[0].text


def test_pptx_parser_builds_slide_parent_and_shape_children(sample_pptx: Path) -> None:
    document = parse_file(sample_pptx)

    assert document.metadata.doc_type == "pptx"
    slides = [block for block in document.blocks if block.block_type == BlockType.SLIDE]
    assert len(slides) == 1
    assert slides[0].metadata.get("slide_title") == "Quarterly Update"

    children = document.children_of(slides[0].block_id)
    assert children
    assert any(child.block_type == BlockType.SHAPE for child in children)

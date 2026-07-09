"""PowerPoint presentation parser."""

from __future__ import annotations

from pathlib import Path

from custom_rag.core.exceptions import ParseFailedError
from custom_rag.core.types import BlockLocation, BlockType, DocumentMetadata, ParsedDocument
from custom_rag.ingestion.parsers.base import BaseParser, BlockBuilder


class PPTXParser(BaseParser):
    name = "pptx"
    supported_extensions = frozenset({".pptx"})
    supported_mimes = frozenset(
        {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }
    )

    def parse(self, path: Path, metadata: DocumentMetadata) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParseFailedError(
                str(path),
                "python-pptx is not installed; use pip install 'custom-rag[ingestion]'",
                cause=exc,
            ) from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ParseFailedError(str(path), "unable to open PPTX file", cause=exc) from exc

        builder = BlockBuilder()
        raw_parts: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_id = f"slide_{slide_index}"
            slide_parts: list[str] = []
            shape_order = 0

            title = _slide_title(slide)
            if title:
                slide_parts.append(title)

            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for row_index, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if not cells:
                                continue
                            row_text = " | ".join(cells)
                            shape_order += 1
                            builder.add(
                                block_id=f"{slide_id}_row_{row_index}",
                                block_type=BlockType.TABLE_ROW,
                                text=row_text,
                                parent_block_id=slide_id,
                                order=shape_order,
                                hierarchy_path=["doc", slide_id, f"row_{row_index}"],
                                location=BlockLocation(slide_number=slide_index),
                                metadata={"row_index": row_index},
                            )
                            slide_parts.append(row_text)
                    continue

                text = shape.text.strip()
                if not text or (title and text == title):
                    continue

                shape_order += 1
                builder.add(
                    block_id=f"{slide_id}_shape_{shape_order}",
                    block_type=BlockType.SHAPE,
                    text=text,
                    parent_block_id=slide_id,
                    order=shape_order,
                    hierarchy_path=["doc", slide_id, f"shape_{shape_order}"],
                    location=BlockLocation(slide_number=slide_index),
                )
                slide_parts.append(text)

            notes_text = _speaker_notes(slide)
            if notes_text:
                shape_order += 1
                builder.add(
                    block_id=f"{slide_id}_notes",
                    block_type=BlockType.NOTE,
                    text=notes_text,
                    parent_block_id=slide_id,
                    order=shape_order,
                    hierarchy_path=["doc", slide_id, "notes"],
                    location=BlockLocation(slide_number=slide_index),
                )
                slide_parts.append(notes_text)

            slide_text = "\n".join(slide_parts)
            builder.add(
                block_id=slide_id,
                block_type=BlockType.SLIDE,
                text=slide_text,
                order=slide_index,
                hierarchy_path=["doc", slide_id],
                location=BlockLocation(slide_number=slide_index),
                metadata={"slide_title": title, "slide_number": slide_index},
            )
            raw_parts.append(slide_text)

        doc_metadata = metadata.model_copy(update={"doc_type": "pptx"})
        return ParsedDocument(
            metadata=doc_metadata,
            blocks=builder.blocks,
            raw_text="\n\n".join(raw_parts) if raw_parts else None,
        )


def _slide_title(slide: object) -> str | None:
    shapes = getattr(slide, "shapes", None)
    if shapes is None:
        return None
    title_shape = getattr(shapes, "title", None)
    if title_shape is None or not getattr(title_shape, "text", ""):
        return None
    return title_shape.text.strip()


def _speaker_notes(slide: object) -> str | None:
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None or notes_slide.notes_text_frame is None:
        return None
    text = notes_slide.notes_text_frame.text.strip()  # type: ignore[union-attr]
    return text or None
